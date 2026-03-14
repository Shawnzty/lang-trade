"""Slide rendering helpers."""

from __future__ import annotations

import math
import textwrap
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont

from utils import ensure_dir, estimate_seconds_from_text


@dataclass
class SlideRenderingResult:
    """Rendered deck artifacts."""

    deck_pptx_path: Path
    slide_images_dir: Path
    thumbnails_dir: Path
    slide_timing_hints: list[dict[str, Any]]


class SlideRenderer:
    """Render PNG slides and a PPTX deck from an approved spec."""

    def __init__(self, config: dict[str, Any]) -> None:
        self.config = config

    def render(self, deck_spec: dict[str, Any], output_dir: Path) -> SlideRenderingResult:
        slide_images_dir = ensure_dir(output_dir / "slide_images")
        thumbnails_dir = ensure_dir(output_dir / "thumbnails")
        pptx_path = output_dir / "deck.pptx"
        presentation = self._build_presentation(deck_spec)
        presentation.save(pptx_path)
        timing_hints: list[dict[str, Any]] = []
        for slide in deck_spec.get("slides", []):
            image_path = slide_images_dir / f"slide-{int(slide['slide_number']):02d}.png"
            thumb_path = thumbnails_dir / f"slide-{int(slide['slide_number']):02d}.png"
            self._render_slide_image(slide, image_path)
            self._render_thumbnail(image_path, thumb_path)
            timing_hints.append(
                {
                    "slide_number": int(slide["slide_number"]),
                    "estimated_duration_sec": float(slide.get("estimated_duration_sec") or 0)
                    or estimate_seconds_from_text(str(slide.get("narration_text", ""))),
                    "title": slide.get("title", ""),
                }
            )
        return SlideRenderingResult(
            deck_pptx_path=pptx_path,
            slide_images_dir=slide_images_dir,
            thumbnails_dir=thumbnails_dir,
            slide_timing_hints=timing_hints,
        )

    def _build_presentation(self, deck_spec: dict[str, Any]) -> Presentation:
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        for slide in deck_spec.get("slides", []):
            pptx_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            self._decorate_pptx_slide(pptx_slide, slide)
        if len(presentation.slides) > len(deck_spec.get("slides", [])):
            presentation.slides._sldIdLst.remove(presentation.slides._sldIdLst[0])  # type: ignore[attr-defined]
        return presentation

    def _decorate_pptx_slide(self, pptx_slide: Any, slide: dict[str, Any]) -> None:
        from pptx.util import Inches, Pt

        width = Inches(13.333)
        height = Inches(7.5)
        background = pptx_slide.shapes.add_shape(1, 0, 0, width, height)
        background.fill.solid()
        background.fill.fore_color.rgb = self._rgb(self.config.get("background_start", "#1f3b4d"))
        background.line.fill.background()
        accent = pptx_slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.55))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self._rgb(self.config.get("accent_color", "#f2a65a"))
        accent.line.fill.background()
        title_box = pptx_slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.8), Inches(1.0))
        title_frame = title_box.text_frame
        title_frame.text = str(slide.get("title", "Untitled"))
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(26)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self._rgb(self.config.get("title_color", "#f7f4ea"))
        body_box = pptx_slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(8.6), Inches(4.3))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        bullets = list(slide.get("bullets", []))
        if slide.get("on_slide_text"):
            bullets.insert(0, str(slide["on_slide_text"]))
        for index, bullet in enumerate(bullets[:6]):
            paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self._rgb(self.config.get("body_color", "#fff9ef"))
        notes_box = pptx_slide.shapes.add_textbox(Inches(9.8), Inches(2.0), Inches(2.5), Inches(4.2))
        notes_frame = notes_box.text_frame
        notes_frame.word_wrap = True
        notes_frame.text = str(slide.get("suggested_visual", ""))
        notes_paragraph = notes_frame.paragraphs[0]
        notes_paragraph.font.size = Pt(15)
        notes_paragraph.font.color.rgb = self._rgb(self.config.get("accent_color", "#f2a65a"))

    def _render_slide_image(self, slide: dict[str, Any], output_path: Path) -> None:
        width = int(self.config.get("width", 1920))
        height = int(self.config.get("height", 1080))
        image = Image.new("RGB", (width, height), self._rgb_tuple(self.config.get("background_start", "#1f3b4d")))
        draw = ImageDraw.Draw(image)
        title_font = self._font(60)
        body_font = self._font(34)
        small_font = self._font(28)
        draw.rounded_rectangle(
            (100, 110, width - 100, height - 110),
            radius=32,
            outline=self._rgb_tuple(self.config.get("accent_color", "#f2a65a")),
            width=5,
        )
        draw.rectangle((100, 110, width - 100, 190), fill=self._rgb_tuple(self.config.get("accent_color", "#f2a65a")))
        draw.text((150, 125), str(slide.get("title", "Untitled")), fill=self._rgb_tuple(self.config.get("title_color", "#f7f4ea")), font=title_font)
        y = 250
        bullets = list(slide.get("bullets", []))
        if slide.get("on_slide_text"):
            bullets.insert(0, str(slide["on_slide_text"]))
        for bullet in bullets[:6]:
            for line_index, line in enumerate(textwrap.wrap(str(bullet), width=64) or [str(bullet)]):
                prefix = "• " if line_index == 0 else "  "
                draw.text((160, y), prefix + line, fill=self._rgb_tuple(self.config.get("body_color", "#fff9ef")), font=body_font)
                y += 46
            y += 10
        visual = textwrap.fill(str(slide.get("suggested_visual", "")), width=28)
        draw.text((1320, 280), visual, fill=self._rgb_tuple(self.config.get("accent_color", "#f2a65a")), font=small_font)
        footer = f"Slide {slide.get('slide_number', '?')}"
        draw.text((150, height - 150), footer, fill=self._rgb_tuple(self.config.get("title_color", "#f7f4ea")), font=small_font)
        image.save(output_path)

    def _render_thumbnail(self, source_path: Path, output_path: Path) -> None:
        image = Image.open(source_path)
        image.thumbnail((480, 270))
        image.save(output_path)

    def _font(self, size: int) -> ImageFont.ImageFont:
        font_candidates = [
            "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
            "/System/Library/Fonts/Supplemental/Helvetica.ttc",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        ]
        for candidate in font_candidates:
            path = Path(candidate)
            if path.exists():
                return ImageFont.truetype(str(path), size=size)
        return ImageFont.load_default()

    def _rgb(self, color: str) -> Any:
        from pptx.dml.color import RGBColor

        r, g, b = self._rgb_tuple(color)
        return RGBColor(r, g, b)

    def _rgb_tuple(self, color: str) -> tuple[int, int, int]:
        color = color.lstrip("#")
        if len(color) != 6:
            return (31, 59, 77)
        return tuple(int(color[index : index + 2], 16) for index in (0, 2, 4))
